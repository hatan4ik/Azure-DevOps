from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).parent / "engineering-intelligence-board-deck.pptx"

SLIDES = [
    ("Engineering Intelligence Transformation Program", ["From reactive engineering to supervised autonomous systems", "18–24 month roadmap for knowledge, quality, reliability, and self-healing infrastructure"]),
    ("The Scaling Problem", ["Knowledge fragments as teams and systems grow", "Senior engineers become operational bottlenecks", "Regression risk, MTTR, onboarding time, and architectural drift increase", "Engineering complexity grows faster than headcount"]),
    ("The Strategic Opportunity", ["Create a Private Engineering Intelligence Platform", "Ground AI in code, tickets, architecture decisions, incidents, and telemetry", "Embed intelligence into IDE, pull requests, CI/CD, and operations", "Measure business outcomes—not AI usage"]),
    ("What This Is — and Is Not", ["Not a chatbot deployment", "Not uncontrolled public AI", "Not a replacement for engineers", "It is secure cognitive infrastructure with human-governed autonomy"]),
    ("Target Architecture", ["Entra ID + private networking + policy enforcement", "Azure OpenAI / enterprise model gateway", "Hybrid retrieval using Azure AI Search or pgvector", "RAG orchestrator with citations and RBAC-aware retrieval", "Observability, Azure DevOps, AKS, IaC, and runbook integrations"]),
    ("Phase 1 — Institutional Memory", ["Continuously ingest repositories, work items, ADRs, documentation, and incidents", "Developer Q&A with citations", "IDE integration and PR summarization", "Goal: reduce search time and onboarding friction"]),
    ("Phase 2 — Intelligent Guardrails", ["PR Guardian agent", "Security and architecture policy checks", "IaC anti-pattern detection", "Historical regression matching", "Goal: catch risk before merge"]),
    ("Phase 3 — Incident Intelligence", ["Correlate logs, metrics, traces, Kubernetes events, and deployment changes", "Summarize failures and rank likely root causes", "Recommend rollback or remediation steps", "Goal: materially reduce MTTR and escalation load"]),
    ("Phase 4–5 — Predictive and Self-Healing", ["Deployment risk scoring and adaptive test depth", "Drift and anomaly detection", "Policy-approved runbook execution", "Automated rollback for known-safe scenarios", "Closed-loop verification and documentation"]),
    ("Financial Impact", ["Value comes from reclaimed engineering capacity, fewer incidents, and faster delivery", "Track cost per query, repo, team, workflow, and agent", "Use small models for routine work and premium models only where justified", "Tie investment gates to measurable ROI"]),
    ("Governance and Safety", ["RBAC before retrieval", "Private endpoints and managed identities", "Prompt-injection defenses and secret redaction", "Human approval for high-blast-radius actions", "Full audit trail, rollback path, and kill switch"]),
    ("Strategic Outcome", ["Persistent institutional knowledge", "Faster and safer software delivery", "Lower operational toil", "Measurable deployment risk", "Progressive autonomy with explicit policy boundaries", "Engineering shifts from reactive operations to supervised autonomous systems"]),
]


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.4), Inches(12.0), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.5), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(14)


def add_footer(slide, n):
    box = slide.shapes.add_textbox(Inches(10.6), Inches(7.05), Inches(2.0), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = f"EIP Transformation | {n}"
    p.font.size = Pt(8)
    p.alignment = PP_ALIGN.RIGHT


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, (title, bullets) in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if idx == 1:
            title_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.7), Inches(11.7), Inches(1.4))
            p = title_box.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            sub = slide.shapes.add_textbox(Inches(1.5), Inches(3.3), Inches(10.3), Inches(1.4))
            tf = sub.text_frame
            for i, text in enumerate(bullets):
                sp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                sp.text = text
                sp.font.size = Pt(20)
                sp.alignment = PP_ALIGN.CENTER
                sp.space_after = Pt(10)
        else:
            add_title(slide, title)
            add_bullets(slide, bullets)
        add_footer(slide, idx)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
